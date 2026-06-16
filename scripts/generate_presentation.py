#!/usr/bin/env python3
"""Generate exam-ready PowerPoint — 10–15 min technical presentation."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "public" / "project"
FIG = OUT_DIR / "outputs"
PPT_PATH = OUT_DIR / "Exam_Presentation_Antenna_RUL.pptx"

NAVY = RGBColor(11, 36, 71)
TEAL = RGBColor(14, 124, 123)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(100, 100, 110)
ACCENT = RGBColor(0, 181, 216)


def _blank_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str = "") -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # header bar
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(0.4))
        sp = st.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = TEAL
        sp.font.italic = True
        top = Inches(1.55)
    else:
        top = Inches(1.25)

    body = slide.shapes.add_textbox(Inches(0.55), top, Inches(8.9), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(30, 30, 40)
        para.space_after = Pt(8)
        para.level = 0


def _image_slide(prs: Presentation, title: str, img_path: Path, caption: str) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(24)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = WHITE

    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.6), Inches(1.1), width=Inches(8.8))

    cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(9), Inches(0.5))
    cp = cap.text_frame.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(12)
    cp.font.color.rgb = GRAY
    cp.font.italic = True


def _title_slide(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(8.8), Inches(2))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "DoE-Based Remaining Useful Life Prediction"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "5G/6G Telecom Tower Antenna Digital Twin"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT
    p2.space_before = Pt(12)

    meta = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(8), Inches(2))
    mf = meta.text_frame
    lines = [
        "Modeling, Simulation & Digital Twin",
        "Prof. Dr. Adele Nasti",
        "Srivardhan Varma Mudunuri · Matric. 100001259",
        "Final Exam Presentation · 10–15 minutes",
    ]
    for i, line in enumerate(lines):
        para = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        para.text = line
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(200, 210, 220)
        para.space_after = Pt(4)


def build() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)

    _blank_slide(
        prs,
        "1. Problem Definition",
        [
            "• Telecom operators lose revenue when 5G/6G antenna towers fail unexpectedly",
            "• Traditional maintenance: periodic climber inspections — costly, risky, reactive",
            "• Industrial context: drone 2D→3D photogrammetry pipeline (USA client, ≈€220K)",
            "• Gap: geometry alone is NOT a digital twin — we need predictive as-used prognostics",
            "• Objective: predict Remaining Useful Life (RUL) under wind + temperature",
        ],
        "Investor pitch: reduce outages, defer site visits, plan maintenance proactively",
    )

    _blank_slide(
        prs,
        "2. Digital Twin — Course Definition",
        [
            '• "Digital representation of a physical product throughout its lifecycle"',
            "• NOT a static CAD copy — spans Design → Manufacture → Build → Test → Service",
            "• Our twin stages (AIAA):",
            "   – As-designed: nominal antenna specifications",
            "   – As-built: drone photogrammetry geometry + QC",
            "   – As-used: physics-based RUL model fed by environmental sensors",
            "• Twin does NOT replace experimental testing — it de-risks decisions",
        ],
    )

    _blank_slide(
        prs,
        "3. Approach & Software Workflow",
        [
            "1. Physics-based degradation model (wind fatigue + thermal aging)",
            "2. Full-factorial Design of Experiments — 4×10 = 40 runs",
            "3. Sensitivity analysis + Pareto ranking of factors",
            "4. Response surface + operating envelope recommendation",
            "5. Python implementation (replicates HEEDS MDO workflow)",
            "",
            "Software: Python (NumPy, Pandas, Matplotlib) + HEEDS methodology",
            "Validation path: compare RUL trends with field inspection data (future work)",
        ],
    )

    _blank_slide(
        prs,
        "4. RUL Physics Model — WHY these equations?",
        [
            "RUL (days) = 175,200 / (f_wind × f_temp × 24)",
            "",
            "f_wind(v) = (v/5)^1.6  — drag-fatigue scaling (structural reliability)",
            "f_temp(T) = 2^((T−20)/15)  — Arrhenius thermal aging (electronics)",
            "",
            "WHY physics-based: interpretable, extrapolatable within bounds, calibratable",
            "Verification: sanity checks — (5m/s,20°C)≈20yr; (35m/s,65°C)≈41 days",
        ],
    )

    _blank_slide(
        prs,
        "5. Design of Experiments (DoE)",
        [
            "Factor A: Wind speed — 4 levels [5, 15, 25, 35] m/s",
            "Factor B: Temperature — 10 levels [20…65] °C in 5°C steps",
            "Design: Full factorial → 40 simulation runs",
            "",
            "WHY full factorial: cheap analytical model; captures interaction for surface",
            "Same workflow as HEEDS Example 4 (Coil Spring DoE)",
            "Process → Parameters → Tagging → Study (DOE) → Run → POST",
        ],
    )

    if (FIG / "05_doe_heatmap.png").exists():
        _image_slide(
            prs,
            "6. DoE Results — Design Space Heatmap",
            FIG / "05_doe_heatmap.png",
            "40-run factorial: RUL (years) across wind × temperature",
        )

    if (FIG / "01_rul_curves.png").exists():
        _image_slide(
            prs,
            "7. RUL vs Temperature (by Wind Level)",
            FIG / "01_rul_curves.png",
            "Higher wind collapses RUL curves — dominant degradation driver",
        )

    if (FIG / "03_pareto_sensitivity.png").exists():
        _image_slide(
            prs,
            "8. Sensitivity & Pareto — Key Finding",
            FIG / "03_pareto_sensitivity.png",
            "Wind speed impact ≈ 3.5× temperature — counterintuitive for RF engineers",
        )

    if (FIG / "02_response_surface_3d.png").exists():
        _image_slide(
            prs,
            "9. Response Surface (DoE POST)",
            FIG / "02_response_surface_3d.png",
            "3D surface — equivalent to HEEDS POST visualization",
        )

    if (FIG / "04_operating_envelope.png").exists():
        _image_slide(
            prs,
            "10. Operating Envelope Recommendation",
            FIG / "04_operating_envelope.png",
            "Safe ops: wind ≤ 12 m/s, temp ≤ 30°C → RUL > 10 years",
        )

    _blank_slide(
        prs,
        "11. Verification vs Validation (This Project)",
        [
            "Verification — 'Have I done the maths right?'",
            "   • Sanity checks at anchor points; unit consistency; monotonic trends",
            "",
            "Validation — 'Have I done the right maths?'",
            "   • Compare predicted RUL ranking with field failure reports",
            "   • Future: correlate with Simcenter FEA fatigue at hotspots",
            "",
            "Calibration: tune exponent 1.6 and Arrhenius slope using failure data",
            "TRL 4–5: component validated in lab/relevant environment; not fleet-deployed",
        ],
    )

    _blank_slide(
        prs,
        "12. HEEDS Mapping & Live Demo",
        [
            "HEEDS Example 4 (Coil Spring) → Antenna RUL:",
            "   • Input tags: wind_ms, temp_C  (like coil_diam, wire_diam)",
            "   • Output tags: rul_days, rul_years  (like deflection, stress)",
            "   • Study type: DOE — NOT optimization",
            "",
            "Live demo: python project/run_project.py",
            "   → prints verification, DoE table, Pareto ratio, exports figures",
            "Jupyter notebook: DigitalTwin_Antenna_RUL.ipynb (full reproducibility)",
        ],
    )

    _blank_slide(
        prs,
        "13. Value Proposition & Creativity",
        [
            "• Extends real €220K drone inspection into predictive maintenance twin",
            "• Creative contribution: telecom-specific RUL physics + DoE envelope",
            "• Pugh matrix winner vs climber-only / geometry-only approaches",
            "• FMEA: wind fatigue failure mode → digital twin monitoring layer",
            "• Business impact: fewer outages, safer crews, data-driven maintenance windows",
        ],
    )

    _blank_slide(
        prs,
        "14. References",
        [
            "1. AIAA (2020). Digital Twin: Definition & Value",
            "2. Grieves & Vickers (2017). Digital Twin: Mitigating Unpredictable Behavior",
            "3. Tao et al. (2019). Digital Twin in Industry. IEEE TII",
            "4. Montgomery (2017). Design and Analysis of Experiments. Wiley",
            "5. Jardine et al. (2006). Machinery diagnostics review. MSSP",
            "6. Liu et al. (2021). UAV photogrammetry for telecom towers",
            "7. Kapteyn et al. (2021). Predictive digital twins. Nature Comp Sci",
            "8. ISO/IEC 30173:2023 — Digital twin terminology",
            "9. Siemens (2024). HEEDS Multi-Disciplinary Design Exploration",
        ],
    )

    _blank_slide(
        prs,
        "15. Thank You — Q&A Backup",
        [
            "Memorise for oral:",
            "• DT definition with 'throughout its lifecycle'",
            "• V vs V vs calibration — one example each from antenna",
            "• Model coupling: structural + thermal + wind + RUL interdependent",
            "• Extrapolation risk beyond 35 m/s or 65°C",
            "• PCA 6 steps · P-diagram centre: maintain structural integrity",
            "",
            "Thank you — questions welcome.",
        ],
    )

    prs.save(str(PPT_PATH))
    return PPT_PATH


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path}")
