#!/usr/bin/env python3
"""
SRH exam presentation — clean layout, readable spacing, official colour palette.
Orange #DF4707 · Navy #0B1D35 · Cream accent · white content slides.
"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "public" / "project" / "ppt-assets"
FIG = ROOT / "public" / "project" / "outputs"
OUT = ROOT / "public" / "project"
PPT_PATH = OUT / "Exam_Presentation_Antenna_RUL_SRH.pptx"
DESKTOP = Path.home() / "Desktop" / "Exam_Presentation_Antenna_RUL_SRH_FINAL.pptx"

# SRH brand (Exercises PDF)
ORANGE = RGBColor(223, 71, 7)
NAVY = RGBColor(11, 29, 53)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(45, 45, 45)
MUTED = RGBColor(90, 90, 90)
CREAM = RGBColor(253, 217, 201)
LIGHT = RGBColor(245, 247, 250)

FONT = "Calibri"
FOOTER = "SRH University  ·  Prof. Dr. Adele Nasti"
STUDENT = "Srivardhan Varma Mudunuri  ·  Matric. 100001259"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ML = Inches(0.75)
MR = Inches(0.75)
HEADER_H = Inches(1.2)
CONTENT_TOP = Inches(1.45)
FOOTER_Y = Inches(7.08)
CONTENT_BOTTOM = Inches(6.85)

LOGO_ORANGE = ASSETS / "srh_logo_orange.png"
LOGO_WHITE = ASSETS / "srh_logo_white.png"


def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, fill: RGBColor, line: bool = False):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if not line:
        s.line.fill.background()
    return s


def _footer(slide, num: int) -> None:
    tb = slide.shapes.add_textbox(ML, FOOTER_Y, Inches(10), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = FOOTER
    p.font.size = Pt(9)
    p.font.name = FONT
    p.font.color.rgb = MUTED

    nb = slide.shapes.add_textbox(Inches(12.1), FOOTER_Y, Inches(0.6), Inches(0.3))
    np = nb.text_frame.paragraphs[0]
    np.text = str(num)
    np.font.size = Pt(9)
    np.font.name = FONT
    np.font.color.rgb = MUTED
    np.alignment = PP_ALIGN.RIGHT


def _logo_orange(slide, top=Inches(0.28)) -> None:
    if LOGO_ORANGE.exists():
        slide.shapes.add_picture(str(LOGO_ORANGE), Inches(11.85), top, height=Inches(0.55))


def _header(slide, title: str, subtitle: str = "") -> None:
    _rect(slide, 0, 0, SLIDE_W, HEADER_H, ORANGE)
    _logo_orange(slide)

    tb = slide.shapes.add_textbox(ML, Inches(0.22), Inches(10.5), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE

    if subtitle:
        sb = slide.shapes.add_textbox(ML, Inches(0.78), Inches(10.5), Inches(0.35))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(15)
        sp.font.name = FONT
        sp.font.color.rgb = CREAM


def _white_bg(slide) -> None:
    _rect(slide, 0, HEADER_H, SLIDE_W, SLIDE_H - HEADER_H, WHITE)


def _bullets(
    slide,
    items: list[str],
    top=CONTENT_TOP,
    width=None,
    font_size=20,
    space_after=14,
) -> None:
    w = width or (SLIDE_W - ML - MR)
    tb = slide.shapes.add_textbox(ML, top, w, CONTENT_BOTTOM - top)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, line in enumerate(items):
        if not line:
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = ""
            para.space_after = Pt(6)
            continue
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(font_size)
        para.font.name = FONT
        para.font.color.rgb = TEXT
        para.space_after = Pt(space_after)
        para.line_spacing = 1.25


def _fit_image(slide, img_path: Path, top, max_w, max_h, left=None):
    if not img_path.exists():
        return
    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = iw / ih
    w = max_w
    h = w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    l = left if left is not None else ML + (max_w - w) / 2
    slide.shapes.add_picture(str(img_path), l, top, width=w, height=h)


def _slide_bullets(prs, num: int, title: str, subtitle: str, bullets: list[str]) -> None:
    slide = _blank(prs)
    _header(slide, title, subtitle)
    _white_bg(slide)
    _bullets(slide, bullets)
    _footer(slide, num)


def _slide_chart(prs, num: int, title: str, subtitle: str, img: Path, caption: str = "") -> None:
    slide = _blank(prs)
    _header(slide, title, subtitle)
    _white_bg(slide)

    max_w = SLIDE_W - ML - MR
    cap_h = Inches(0.45) if caption else Inches(0)
    max_h = CONTENT_BOTTOM - CONTENT_TOP - cap_h - Inches(0.15)
    _fit_image(slide, img, CONTENT_TOP + Inches(0.1), max_w, max_h)

    if caption:
        cb = slide.shapes.add_textbox(ML, CONTENT_BOTTOM - Inches(0.42), max_w, Inches(0.4))
        cp = cb.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(13)
        cp.font.italic = True
        cp.font.color.rgb = MUTED
        cp.font.name = FONT

    _footer(slide, num)


def _slide_split(
    prs,
    num: int,
    title: str,
    left_bullets: list[str],
    img: Path,
) -> None:
    slide = _blank(prs)
    _header(slide, title)
    _white_bg(slide)

    col_w = Inches(5.6)
    _bullets(slide, left_bullets, top=CONTENT_TOP, width=col_w, font_size=18, space_after=12)

    max_w = Inches(6.2)
    max_h = CONTENT_BOTTOM - CONTENT_TOP
    _fit_image(slide, img, CONTENT_TOP, max_w, max_h, left=Inches(6.85))

    _footer(slide, num)


def _title_slide(prs) -> None:
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _rect(slide, 0, Inches(6.55), SLIDE_W, Inches(0.95), ORANGE)

    if LOGO_WHITE.exists():
        slide.shapes.add_picture(str(LOGO_WHITE), ML, Inches(0.45), height=Inches(0.7))

    tb = slide.shapes.add_textbox(ML, Inches(1.85), Inches(11), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DoE-Based Remaining Useful Life Prediction"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "5G/6G Telecom Tower Antenna Digital Twin"
    p2.font.size = Pt(24)
    p2.font.name = FONT
    p2.font.color.rgb = CREAM
    p2.space_before = Pt(14)

    meta = slide.shapes.add_textbox(ML, Inches(4.35), Inches(9), Inches(1.8))
    mf = meta.text_frame
    for i, line in enumerate(
        [
            "Modeling, Simulation & Digital Twin",
            "Prof. Dr. Adele Nasti",
            STUDENT,
            "Final Exam · 10–15 minutes",
        ]
    ):
        para = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        para.text = line
        para.font.size = Pt(15)
        para.font.name = FONT
        para.font.color.rgb = RGBColor(210, 220, 230)
        para.space_after = Pt(6)

    sb = slide.shapes.add_textbox(ML, Inches(6.72), Inches(11), Inches(0.5))
    sp = sb.text_frame.paragraphs[0]
    sp.text = STUDENT
    sp.font.size = Pt(13)
    sp.font.name = FONT
    sp.font.color.rgb = WHITE


def _thankyou(prs) -> None:
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _rect(slide, 0, 0, SLIDE_W, Inches(0.12), ORANGE)

    if LOGO_WHITE.exists():
        slide.shapes.add_picture(str(LOGO_WHITE), ML, Inches(0.45), height=Inches(0.65))

    tb = slide.shapes.add_textbox(ML, Inches(2.8), Inches(10), Inches(1.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = FONT
    p.font.color.rgb = WHITE

    p2 = tb.text_frame.add_paragraph()
    p2.text = "Questions welcome"
    p2.font.size = Pt(22)
    p2.font.name = FONT
    p2.font.color.rgb = CREAM
    p2.space_before = Pt(12)


def build() -> Path:
    assets_script = ROOT / "scripts" / "generate_srh_assets.py"
    run_script = ROOT / "project" / "run_project.py"
    if assets_script.exists():
        subprocess.run([sys.executable, str(assets_script)], check=True)
    if run_script.exists():
        subprocess.run([sys.executable, str(run_script)], check=True)

    prs = _prs()
    n = 1

    _title_slide(prs)
    n += 1

    _slide_bullets(
        prs,
        n,
        "1 · Problem",
        "Predictive maintenance for 5G/6G telecom towers",
        [
            "Tower failures cause outages and revenue loss",
            "Climber inspections are costly, slow, and risky",
            "Drone 2D→3D scan gives geometry — but geometry alone is NOT a digital twin",
            "Goal: predict Remaining Useful Life (RUL) from wind + temperature",
        ],
    )
    n += 1

    _slide_split(
        prs,
        n,
        "2 · Digital Twin Definition",
        [
            "Digital representation throughout its lifecycle",
            "",
            "NOT a static CAD copy",
            "",
            "Stages: Design → Manufacture → Build → Test → Service",
            "",
            "Twin does NOT replace experimental testing",
        ],
        ASSETS / "viz_lifecycle.png",
    )
    n += 1

    _slide_chart(
        prs,
        n,
        "3 · End-to-End Approach",
        "Drone geometry + physics RUL + Design of Experiments",
        ASSETS / "viz_workflow.png",
        "Pipeline extends real industrial drone inspection into predictive maintenance",
    )
    n += 1

    _slide_bullets(
        prs,
        n,
        "4 · Physics-Based RUL Model",
        "Wind fatigue × thermal aging",
        [
            "RUL (days) = 175,200 / (f_wind × f_temp × 24)",
            "",
            "f_wind(v) = (v / 5)^1.6        f_temp(T) = 2^((T − 20) / 15)",
            "",
            "Verification sanity checks:",
            "   •  5 m/s, 20 °C  →  ~20 years",
            "   • 15 m/s, 35 °C  →  ~1.7 years",
            "   • 35 m/s, 65 °C  →  ~41 days",
        ],
    )
    n += 1

    _slide_bullets(
        prs,
        n,
        "5 · Design of Experiments",
        "Full factorial — same workflow as HEEDS Example 5 (Coil Spring)",
        [
            "Factor A: Wind speed — 4 levels [5, 15, 25, 35] m/s",
            "Factor B: Temperature — 10 levels [20…65] °C",
            "Design: 4 × 10 = 40 simulation runs",
            "",
            "HEEDS workflow: Process → Parameters → Tagging → Study → Run → POST",
        ],
    )
    n += 1

    _slide_chart(
        prs,
        n,
        "6 · DoE Results — Design Space",
        "40-run factorial heatmap",
        FIG / "05_doe_heatmap.png",
        "RUL (years) across wind × temperature",
    )
    n += 1

    _slide_chart(
        prs,
        n,
        "7 · RUL vs Temperature",
        "Separate curve for each wind level",
        FIG / "01_rul_curves.png",
        "Higher wind collapses remaining life — dominant degradation driver",
    )
    n += 1

    _slide_chart(
        prs,
        n,
        "8 · Sensitivity Analysis",
        "Pareto ranking of factors",
        FIG / "03_pareto_sensitivity.png",
        "Wind impact ≈ 3.5× temperature",
    )
    n += 1

    _slide_chart(
        prs,
        n,
        "9 · Operating Envelope",
        "Safe operating recommendation",
        FIG / "04_operating_envelope.png",
        "Safe ops: wind ≤ 12 m/s and temp ≤ 30 °C → RUL > 10 years",
    )
    n += 1

    _slide_chart(
        prs,
        n,
        "10 · HEEDS Mapping",
        "Coil Spring DoE ≡ Antenna factorial study",
        ASSETS / "viz_heeds_map.png",
    )
    n += 1

    _slide_chart(
        prs,
        n,
        "11 · Key Results",
        "Predictive maintenance value · TRL 4–5",
        ASSETS / "viz_key_results.png",
    )
    n += 1

    _slide_split(
        prs,
        n,
        "12 · Verification & Validation",
        [
            "Verification — maths right?",
            "Sanity checks at anchor points",
            "",
            "Validation — right maths?",
            "Compare RUL vs field failures",
            "",
            "Calibration ≠ validation",
        ],
        ASSETS / "viz_vv.png",
    )
    n += 1

    _slide_bullets(
        prs,
        n,
        "13 · Live Demo & Value",
        "Reproducible Python pipeline",
        [
            "Run: python project/run_project.py",
            "Notebook: DigitalTwin_Antenna_RUL.ipynb",
            "",
            "Value: fewer outages, safer crews, data-driven maintenance",
            "Creative: telecom-specific RUL physics + DoE envelope",
        ],
    )
    n += 1

    _slide_bullets(
        prs,
        n,
        "14 · References",
        "",
        [
            "1. AIAA (2020). Digital Twin: Definition & Value",
            "2. Grieves & Vickers (2017). Digital Twin: Mitigating Unpredictable Behavior",
            "3. Tao et al. (2019). Digital Twin in Industry. IEEE TII",
            "4. Montgomery (2017). Design and Analysis of Experiments",
            "5. Liu et al. (2021). UAV photogrammetry for telecom towers",
            "6. Kapteyn et al. (2021). Predictive digital twins. Nature Comp Sci",
            "7. Siemens (2024). HEEDS Multi-Disciplinary Design Exploration",
        ],
    )
    n += 1

    _thankyou(prs)

    OUT.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPT_PATH))
    prs.save(str(DESKTOP))
    return PPT_PATH


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path}")
    print(f"Desktop: {DESKTOP}")
